from pptx import Presentation
from collections.abc import Generator
import os


class PowerpointParser:
    """
    This class is used to parse a powerpoint file and extract the text from each slide
    """
    pptx_file: str

    def __init__(self, file_path: str) -> None:
        """
        Initialize the class
        :param file_path: The path to the PowerPoint file
        """

        if os.path.isfile(file_path):
            self.pptx_file = file_path
        else:
            raise FileNotFoundError("The powerpoint file was not found")

    def extract_text_from_slide(self) -> Generator:
        """
        Extract the text from each slide
        :return: A generator that yields the text from each slide
        """

        # Open the PowerPoint file
        with open(self.pptx_file, "rb") as presentation_file:
            presentation = Presentation(presentation_file)

            # Iterate through each slide and extract the text
            for slide in presentation.slides:
                slide_text = ""
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                slide_text += '\n' + run.text.strip()

                if slide_text != "":
                    yield slide_text


def main():
    path = input("Enter your powerpoint file path\n")

    pptx_object = PowerpointParser(path)

    for slide, text in enumerate(pptx_object.extract_text_from_slide()):
        print(f"[slide {slide}]{text}")


if __name__ == "__main__":
    main()
